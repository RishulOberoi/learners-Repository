from flask import Flask, request, jsonify, send_file, render_template
import pandas as pd
import re
import io
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image
import nltk
import os
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Download required NLTK data
try:
    nltk.download('punkt', quiet=True)
except:
    pass

from transformers import pipeline

# Global variables for models
summarizer = None
expander = None

def load_models():
    global summarizer, expander
    try:
        summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
        expander = pipeline("text2text-generation", model="t5-base")
    except Exception as e:
        print(f"Error loading models: {e}")

def remove_unwanted_phrases(text):
    unwanted_phrases = [
        "For confidential support call the Samaritans on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org",
        "For confidential support, call the Samaritans on 08457 90 90, visit a local Samaritans branch or see www.samaritans.org",
        "For confidential support call the Samaritans in the UK on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org for details.",
        "For confidential support, call the Samaritans in the UK on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org for details.",
        "For confidential support call the Samaritans",
        "For confidential support, call the Samaritans",
        "visit a local Samaritans branch or see www.samaritans.org",
        "visit a local Samaritans branch or see www.samaritans.org for details"
    ]
    for phrase in unwanted_phrases:
        text = text.replace(phrase, "")
    text = re.sub(r'\bFor confidential support.*?samaritans\.org\.?', '', text, flags=re.IGNORECASE)
    text = re.sub(r'contact (us|this) .*?samaritans\.org\.?', '', text, flags=re.IGNORECASE)
    return text.strip()

def summarize_text(text):
    global summarizer
    try:
        if summarizer is None:
            return remove_unwanted_phrases(text)
        summary = summarizer(
            text, max_length=300, min_length=100, length_penalty=0.9,
            no_repeat_ngram_size=3, do_sample=False
        )
        result = summary[0]['summary_text'].strip()
        return remove_unwanted_phrases(result)
    except Exception as e:
        print(f"Summarization failed: {e}")
        return remove_unwanted_phrases(text)

def expand_text(text):
    global expander
    prompt = f"Provide a more detailed explanation: {text}"
    try:
        if expander is None:
            return remove_unwanted_phrases(text)
        expanded = expander(prompt, max_new_tokens=120, do_sample=False)
        result = expanded[0]['generated_text'].strip()
        if result.lower().startswith(prompt.lower()) or len(result) <= len(text) + 8:
            return remove_unwanted_phrases(text)
        return remove_unwanted_phrases(result)
    except Exception as e:
        print(f"Expansion failed: {e}")
        return remove_unwanted_phrases(text)

def parse_content_to_bullets(text):
    bullets = []
    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue
        if re.match(r'^[-*]\s', stripped):
            bullets.append(('sub', stripped[1:].strip()))
        else:
            bullets.append(('main', stripped))
    return bullets

def add_image_autofit(slide, image_path, left_inches, top_inches, width_inches, height_inches, min_size=1.0):
    try:
        if pd.isna(image_path) or not str(image_path).strip():
            return
        if str(image_path).startswith('http'):
            resp = requests.get(image_path)
            img = Image.open(io.BytesIO(resp.content))
        else:
            img = Image.open(image_path)
        img_w, img_h = img.size
        dpi = img.info.get('dpi', (96, 96))[0] or 96
        img_width_in = img_w / dpi
        img_height_in = img_h / dpi
        box_aspect = width_inches / height_inches
        img_aspect = img_width_in / img_height_in
        if img_aspect > box_aspect:
            final_width = max(min(width_inches, img_width_in), min_size)
            final_height = final_width / img_aspect
            if final_height < min_size:
                final_height = min_size
                final_width = final_height * img_aspect
        else:
            final_height = max(min(height_inches, img_height_in), min_size)
            final_width = final_height * img_aspect
            if final_width < min_size:
                final_width = min_size
                final_height = final_width / img_aspect
        pos_left = left_inches + max((width_inches - final_width) / 2, 0)
        pos_top = top_inches + max((height_inches - final_height) / 2, 0)
        img_obj = io.BytesIO(resp.content) if str(image_path).startswith('http') else image_path
        slide.shapes.add_picture(
            img_obj,
            Inches(pos_left), Inches(pos_top),
            width=Inches(final_width), height=Inches(final_height)
        )
    except Exception as e:
        print(f"Could not add image {image_path}: {e}")

def create_slide(prs, title, content_bullets, image_path=None):
    if not isinstance(title, str) or title.strip() == "" or pd.isna(title):
        title = "Untitled Slide"
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    para = title_shape.text_frame.paragraphs[0]
    para.font.size = Pt(34)
    para.font.bold = True
    para.font.color.rgb = RGBColor(10, 60, 120)
    slide_width_in = prs.slide_width / 914400
    slide_height_in = prs.slide_height / 914400
    margin = 0.5
    gap = 0.3
    MAX_BULLETS = 10
    if image_path and str(image_path).strip():
        img_area_width = max(slide_width_in * 0.38, 2.5)
        text_area_width = slide_width_in - img_area_width - gap - 2 * margin
        text_area_left = margin
        text_area_top = 1.5
        text_area_height = slide_height_in - text_area_top - margin
        img_left = text_area_left + text_area_width + gap
        img_top = text_area_top
        img_width = img_area_width
        img_height = text_area_height
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                slide.shapes._spTree.remove(shape._element)
                break
        text_box = slide.shapes.add_textbox(
            Inches(text_area_left), Inches(text_area_top),
            Inches(text_area_width), Inches(text_area_height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        for i, (level, text) in enumerate(content_bullets):
            if i >= MAX_BULLETS:
                paragraph = text_frame.add_paragraph()
                paragraph.text = "... (content truncated)"
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(150, 150, 150)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                break
            paragraph = text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0 if level == 'main' else 1
            paragraph.font.size = Pt(20) if level == 'main' else Pt(18)
            paragraph.font.color.rgb = RGBColor(35, 35, 35)
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        add_image_autofit(slide, image_path, img_left, img_top, img_width, img_height)
    else:
        text_area_left = margin
        text_area_top = 1.5
        text_area_width = slide_width_in - 2 * margin
        text_area_height = slide_height_in - text_area_top - margin
        text_box = slide.shapes.add_textbox(
            Inches(text_area_left), Inches(text_area_top),
            Inches(text_area_width), Inches(text_area_height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        for i, (level, text) in enumerate(content_bullets):
            if i >= MAX_BULLETS:
                paragraph = text_frame.add_paragraph()
                paragraph.text = "... (content truncated)"
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(150, 150, 150)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                break
            paragraph = text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0 if level == 'main' else 1
            paragraph.font.size = Pt(20) if level == 'main' else Pt(18)
            paragraph.font.color.rgb = RGBColor(35, 35, 35)
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    return slide

def decide_enrichment(title, content_raw):
    if not content_raw or str(content_raw).lower() == "nan":
        return "Content not available."
    cleaned_text = remove_unwanted_phrases(content_raw)
    length = len(cleaned_text)
    if length < 80:
        print(f"[{title}] Expanding short content ({length} chars).")
        return expand_text(cleaned_text)
    if length > 100:
        print(f"[{title}] Summarizing long content ({length} chars).")
        return summarize_text(cleaned_text)
    print(f"[{title}] Using as-is ({length} chars).")
    return cleaned_text

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if file.filename.endswith('.csv'):
            try:
                df = pd.read_csv(file, encoding='utf-8')
            except UnicodeDecodeError:
                file.seek(0)
                try:
                    df = pd.read_csv(file, encoding='utf-16')
                except UnicodeDecodeError:
                    file.seek(0)
                    df = pd.read_csv(file, encoding='latin1')
        elif file.filename.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(file)
        else:
            return jsonify({'error': 'Invalid file format. Please upload CSV or Excel file.'}), 400
        
        df.columns = df.columns.str.strip()
        
        # Convert DataFrame to HTML for preview
        table_html = df.head(10).to_html(classes='table table-striped', table_id='dataPreview')
        
        return jsonify({
            'success': True,
            'table_html': table_html,
            'columns': list(df.columns),
            'row_count': len(df)
        })
    
    except Exception as e:
        return jsonify({'error': f'Error processing file: {str(e)}'}), 500

@app.route('/generate', methods=['POST'])
def generate_presentation():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename.endswith('.csv'):
            try:
                df = pd.read_csv(file, encoding='utf-8')
            except UnicodeDecodeError:
                file.seek(0)
                try:
                    df = pd.read_csv(file, encoding='utf-16')
                except UnicodeDecodeError:
                    file.seek(0)
                    df = pd.read_csv(file, encoding='latin1')
        else:
            df = pd.read_excel(file)
        
        df.columns = df.columns.str.strip()
        
        prs = Presentation()
        cols_lower = set(c.lower() for c in df.columns)
        
        if 'title' in cols_lower and 'content' in cols_lower:
            for idx, row in df.iterrows():
                title_raw = row.get('Title') or row.get('title') or 'Untitled Slide'
                title = str(title_raw) if not pd.isna(title_raw) else 'Untitled Slide'

                content_raw_val = row.get('Content') or row.get('content') or ''
                content_raw = str(content_raw_val) if not pd.isna(content_raw_val) else ''

                content_used = decide_enrichment(title, content_raw)
                bullets = parse_content_to_bullets(content_used)

                image_path_val = row.get('Image') or row.get('image') or None
                image_path = image_path_val if (image_path_val and not pd.isna(image_path_val)) else None

                create_slide(prs, title, bullets, image_path)
        else:
            # For generic tabular data, create a single table slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title-only layout
            slide.shapes.title.text = "Table View"
            rows, cols = df.shape
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
            # Add headers
            for c, name in enumerate(df.columns):
                table.cell(0, c).text = str(name)
            # Add data
            for r in range(rows):
                for c in range(cols):
                    table.cell(r + 1, c).text = str(df.iat[r, c])

        # Save presentation
        pptx_file = "PowerPoint_Generator.pptx"
        prs.save(pptx_file)
        
        return send_file(
            pptx_file,
            as_attachment=True,
            download_name="PowerPoint_Generator.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    except Exception as e:
        return jsonify({'error': f'Error generating presentation: {str(e)}'}), 500

if __name__ == '__main__':
    print("Loading AI models... This may take a few minutes on first run.")
    load_models()
    print("✅ Models loaded successfully!")
    app.run(debug=True, port=5000)
