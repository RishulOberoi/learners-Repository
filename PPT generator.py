import pandas as pd
import re
import io
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from tkinter import Tk, filedialog
from PIL import Image
import nltk

nltk.download('punkt')

from transformers import pipeline

# Load transformer pipelines once
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
expander = pipeline("text2text-generation", model="t5-base")

def remove_unwanted_phrases(text):
    unwanted_phrases = [
        "For confidential support call the Samaritans on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org.",
        "For confidential support call the Samaritans on 08457 90 90, visit a local Samaritans branch or see www.samaritans.org.",
        "For confidential support, call the Samaritans on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org.",
        "For confidential support, call the Samaritans on 08457 90 90, visit a local Samaritans branch or see www.samaritans.org.",
        "For confidential support call the Samaritans in the UK on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org",
        "For confidential support call the Samaritans in the UK on 08457 90 90 90, visit a local Samaritans branch or see www.samaritans.org.",
        # Add any other variants if needed
    ]
    for phrase in unwanted_phrases:
        text = text.replace(phrase, "")
    return text.strip()

def summarize_text(text):
    try:
        summary = summarizer(
            text,
            max_length=300,
            min_length=100,
            length_penalty=0.9,
            no_repeat_ngram_size=3,
            do_sample=False
        )
        result = summary[0]['summary_text'].strip()
        result = remove_unwanted_phrases(result)
        return result
    except Exception as e:
        print(f"Summarization failed: {e}")
        return remove_unwanted_phrases(text)

def expand_text(text):
    prompt = f"Provide a more detailed explanation: {text}"
    try:
        expanded = expander(prompt, max_new_tokens=120, do_sample=False)
        result = expanded[0]['generated_text'].strip()
        # Fallback if expansion repeats prompt or is too close to input
        if result.lower().startswith(prompt.lower()) or len(result) <= len(text) + 8:
            print(f"[EXPAND FAILED for text: {text[:30]}..., using original]")
            return remove_unwanted_phrases(text)
        return remove_unwanted_phrases(result)
    except Exception as e:
        print(f"Expansion failed: {e}")
        return remove_unwanted_phrases(text)

def parse_content_to_bullets(text):
    bullets = []
    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue
        if re.match(r'^[-*]\s', stripped):
            bullets.append(('sub', stripped[1:].strip()))
        else:
            bullets.append(('main', stripped))
    return bullets

def add_image_autofit(slide, image_path, left_inches, top_inches, width_inches, height_inches, min_size=1.0):
    try:
        if pd.isna(image_path) or not str(image_path).strip():
            return
        if str(image_path).startswith('http'):
            resp = requests.get(image_path)
            img = Image.open(io.BytesIO(resp.content))
        else:
            img = Image.open(image_path)
        img_w, img_h = img.size
        dpi = img.info.get('dpi', (96, 96))[0] or 96
        img_width_in = img_w / dpi
        img_height_in = img_h / dpi
        box_aspect = width_inches / height_inches
        img_aspect = img_width_in / img_height_in

        if img_aspect > box_aspect:
            final_width = max(min(width_inches, img_width_in), min_size)
            final_height = final_width / img_aspect
            if final_height < min_size:
                final_height = min_size
                final_width = final_height * img_aspect
        else:
            final_height = max(min(height_inches, img_height_in), min_size)
            final_width = final_height * img_aspect
            if final_width < min_size:
                final_width = min_size
                final_height = final_width / img_aspect

        pos_left = left_inches + max((width_inches - final_width) / 2, 0)
        pos_top = top_inches + max((height_inches - final_height) / 2, 0)

        pic_image = io.BytesIO(resp.content) if str(image_path).startswith('http') else image_path

        slide.shapes.add_picture(
            pic_image,
            Inches(pos_left), Inches(pos_top),
            width=Inches(final_width), height=Inches(final_height)
        )
    except Exception as e:
        print(f"Could not add image {image_path}: {e}")

def create_slide(prs, title, content_bullets, image_path=None):
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    para = title_shape.text_frame.paragraphs[0]
    para.font.size = Pt(34)
    para.font.bold = True
    para.font.color.rgb = RGBColor(10, 60, 120)

    slide_width_in = prs.slide_width / 914400
    slide_height_in = prs.slide_height / 914400
    margin = 0.5
    gap = 0.3
    MAX_BULLETS = 10  # limit to prevent overflow

    if image_path and str(image_path).strip():
        img_area_width = max(slide_width_in * 0.38, 2.5)
        text_area_width = slide_width_in - img_area_width - gap - 2 * margin
        text_area_left = margin
        text_area_top = 1.5
        text_area_height = slide_height_in - text_area_top - margin

        img_left = text_area_left + text_area_width + gap
        img_top = text_area_top
        img_width = img_area_width
        img_height = text_area_height

        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                slide.shapes._spTree.remove(shape._element)
                break

        text_box = slide.shapes.add_textbox(
            Inches(text_area_left), Inches(text_area_top),
            Inches(text_area_width), Inches(text_area_height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        for i, (level, text) in enumerate(content_bullets):
            if i >= MAX_BULLETS:
                paragraph = text_frame.add_paragraph()
                paragraph.text = "... (content truncated)"
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(150, 150, 150)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                break
            paragraph = text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0 if level == 'main' else 1
            paragraph.font.size = Pt(20) if level == 'main' else Pt(18)
            paragraph.font.color.rgb = RGBColor(35, 35, 35)
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        add_image_autofit(slide, image_path, img_left, img_top, img_width, img_height)
    else:
        text_area_left = margin
        text_area_top = 1.5
        text_area_width = slide_width_in - 2 * margin
        text_area_height = slide_height_in - text_area_top - margin

        text_box = slide.shapes.add_textbox(
            Inches(text_area_left), Inches(text_area_top),
            Inches(text_area_width), Inches(text_area_height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        for i, (level, text) in enumerate(content_bullets):
            if i >= MAX_BULLETS:
                paragraph = text_frame.add_paragraph()
                paragraph.text = "... (content truncated)"
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(150, 150, 150)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                break
            paragraph = text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0 if level == 'main' else 1
            paragraph.font.size = Pt(20) if level == 'main' else Pt(18)
            paragraph.font.color.rgb = RGBColor(35, 35, 35)
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    return slide

def ask_for_file():
    root = Tk()
    root.withdraw()
    file_path = filedialog.askopenfilename(
        title='Select Excel or CSV file for PPT generation',
        filetypes=[('CSV or Excel files', '*.csv;*.xlsx'), ('All Files', '*.*')]
    )
    root.destroy()
    return file_path

def decide_enrichment(title, content_raw):
    if not content_raw or str(content_raw).lower() == "nan":
        return "Content not available."
    cleaned_text = remove_unwanted_phrases(content_raw)
    length = len(cleaned_text)
    if length < 80:
        print(f"[{title}] Expanding short content ({length} chars).")
        return expand_text(cleaned_text)
    if length > 100:
        print(f"[{title}] Summarizing long content ({length} chars).")
        return summarize_text(cleaned_text)
    print(f"[{title}] Using as-is ({length} chars).")
    return cleaned_text

def main():
    file_path = ask_for_file()
    if not file_path:
        print("No file selected. Exiting.")
        return

    if file_path.endswith('.csv'):
        try:
            df = pd.read_csv(file_path, encoding='utf-8')
        except UnicodeDecodeError:
            try:
                df = pd.read_csv(file_path, encoding='cp1252')
            except UnicodeDecodeError:
                df = pd.read_csv(file_path, encoding='latin1')
    else:
        df = pd.read_excel(file_path)

    df.columns = df.columns.str.strip()
    prs = Presentation()

    for idx, row in df.iterrows():
        title = str(row.get('Title', 'Untitled Slide'))
        content_raw = str(row.get('Content', ''))
        content_used = decide_enrichment(title, content_raw)
        bullets = parse_content_to_bullets(content_used)
        image_path = row.get('Image', None)
        create_slide(prs, title, bullets, image_path)

    output_file = 'Output PPT.pptx'
    prs.save(output_file)
    print(f'Presentation generated: {output_file}')

if __name__ == '__main__':
    main()
